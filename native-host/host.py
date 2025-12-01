#!/usr/bin/env python3
"""
Native messaging host for UltraThink Chrome extension.

This module handles:
- Native messaging communication with Chrome extension
- Saving entries to kb.md markdown knowledge base
- AI processing pipeline (grammar correction, summarisation, classification)
- File and screenshot management
- Desktop widget lifecycle

The host receives JSON messages from the extension via stdin and sends
responses via stdout. All AI processing runs in background threads to
avoid blocking the UI.

Entry format in kb.md:
    - `type` | `source` | `timestamp` | [Title](URL)
      - Notes: User's notes
      - Entity: project|task|knowledge
      - Topics: Topic1, Topic2
      - AI Summary: Generated summary
"""

import sys
import json
import struct
import os
import re
import base64
import urllib.request
import urllib.error
import urllib.parse
import threading
from pathlib import Path
from datetime import datetime


# =============================================================================
# SECURITY: Path and Input Validation
# =============================================================================

# Blocked system directories (case-insensitive)
BLOCKED_PATHS = [
    'c:\\windows',
    'c:\\program files',
    'c:\\program files (x86)',
    'c:\\programdata',
    '/etc',
    '/usr',
    '/bin',
    '/sbin',
    '/var',
]

# Type mappings for AI Summary strategies
IMAGE_TYPES = ['screenshot', 'image']
RESEARCH_TYPES = ['long-note']
AUDIO_TYPES = ['audio']
DOCUMENT_TYPES = ['pdf', 'markdown', 'ms-word', 'ms-excel', 'ms-powerpoint', 'ms-onenote']
LINK_TYPES = ['link', 'chatgpt', 'claude', 'perplexity', 'notion']  # Types that need URL browsing
TEXT_TYPES = ['snippet', 'note', 'para', 'idea']  # Types that just need text summarisation
NO_SUMMARY_TYPES = ['video']  # Types that don't get AI summary

# Settings file path (in same directory as this script)
SETTINGS_FILE = Path(__file__).parent / 'settings.json'


def load_settings():
    """
    Load settings from settings.json file.

    Returns:
        dict: Settings dictionary, or empty dict if file not found/invalid.

    Settings include:
        - project_folder: Path to kb.md directory (single source of truth)
        - openai_api_key: OpenAI API key for AI processing
        - github_token, notion_token, fastmail_token: Integration tokens
        - classification_prompt: Custom prompt for entry classification
        - grammar_prompt: Custom prompt for grammar correction
        - Other AI prompt customisations
    """
    try:
        if SETTINGS_FILE.exists():
            with open(SETTINGS_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception:
        pass
    return {}


def save_settings(settings):
    """
    Save settings to settings.json file.

    Args:
        settings (dict): Settings dictionary to save.

    Returns:
        bool: True if saved successfully, False otherwise.
    """
    try:
        with open(SETTINGS_FILE, 'w', encoding='utf-8') as f:
            json.dump(settings, f, indent=2)
        return True
    except Exception:
        return False


def update_settings(updates):
    """
    Update specific settings without overwriting others.

    Args:
        updates (dict): Dictionary of settings to update.

    Returns:
        bool: True if saved successfully, False otherwise.
    """
    settings = load_settings()
    settings.update(updates)
    return save_settings(settings)


def get_prompt(settings, key):
    """
    Get a prompt from settings.json.

    Args:
        settings (dict): Settings dictionary from load_settings().
        key (str): Settings key for the prompt (e.g., 'classification_prompt').

    Returns:
        str: The prompt from settings.
    """
    return settings.get(key, '').strip()


def validate_project_folder(project_folder):
    """
    Validate that a project folder path is safe and usable.

    Performs security checks to prevent path traversal attacks and
    writing to system directories.

    Args:
        project_folder (str): Path to validate.

    Returns:
        tuple: (is_valid: bool, result: Path|str)
            If valid, result is the resolved Path object.
            If invalid, result is an error message string.

    Security checks:
        - Must be absolute path
        - Must exist as a directory
        - Cannot be a system directory (Windows, Program Files, /etc, etc.)
        - Cannot contain '..' path traversal
    """
    if not project_folder:
        return False, "Project folder cannot be empty"

    try:
        folder_path = Path(project_folder).resolve()

        # Must be absolute
        if not folder_path.is_absolute():
            return False, "Project folder must be an absolute path"

        # Must exist and be a directory
        if not folder_path.exists():
            return False, f"Project folder does not exist: {project_folder}"

        if not folder_path.is_dir():
            return False, f"Project folder is not a directory: {project_folder}"

        # Block system directories
        folder_str = str(folder_path).lower()
        for blocked in BLOCKED_PATHS:
            if folder_str.startswith(blocked):
                return False, f"Cannot use system directory: {folder_path}"

        # Block path traversal attempts
        if '..' in project_folder:
            return False, "Path cannot contain '..'"

        return True, folder_path

    except Exception as e:
        return False, f"Invalid path: {str(e)}"


def sanitize_filename(filename):
    """
    Sanitize a filename to prevent path traversal and security issues.

    Args:
        filename (str): Original filename to sanitize.

    Returns:
        str: Safe filename with only alphanumeric chars, dashes,
            underscores, dots, and spaces. Max 200 characters.

    Security measures:
        - Strips path components (only keeps basename)
        - Removes path separators and null bytes
        - Removes leading dots (hidden files)
        - Limits length to 200 characters
    """
    if not filename:
        return 'unnamed_file'

    # Get just the basename (strip any path components)
    basename = Path(filename).name

    # Remove any remaining path separators and null bytes
    basename = basename.replace('/', '').replace('\\', '').replace('\x00', '')

    # Remove leading dots (hidden files / parent traversal)
    basename = basename.lstrip('.')

    # Whitelist: only allow alphanumeric, dash, underscore, dot, space
    basename = re.sub(r'[^a-zA-Z0-9\-_\. ]', '_', basename)

    # Limit length
    if len(basename) > 200:
        name_parts = basename.rsplit('.', 1)
        if len(name_parts) == 2:
            basename = name_parts[0][:190] + '.' + name_parts[1][:10]
        else:
            basename = basename[:200]

    # Ensure not empty after sanitization
    if not basename or basename == '.':
        basename = 'unnamed_file'

    return basename


def validate_entry(entry):
    """
    Validate entry object has required fields.
    Returns (is_valid, error_message).
    """
    if not isinstance(entry, dict):
        return False, "Entry must be an object"

    required_fields = ['type', 'captured', 'source']
    for field in required_fields:
        if field not in entry:
            return False, f"Missing required field: {field}"

    # Validate timestamp format (YYYY-MM-DD HH:MM:SS)
    timestamp = entry.get('captured', '')
    if not re.match(r'^\d{4}-\d{2}-\d{2} \d{2}:\d{2}:\d{2}$', timestamp):
        return False, f"Invalid timestamp format: {timestamp}"

    return True, None


# =============================================================================
# External Service Integrations
# =============================================================================

def search_github(query, github_token, repos=None, max_results=10):
    """
    Search GitHub for issues, commits, and code related to a query.

    Args:
        query (str): Search query string.
        github_token (str): GitHub Personal Access Token.
        repos (str|None): Comma-separated list of repos to search (e.g., "owner/repo1, owner/repo2").
                         If None, searches all accessible repos.
        max_results (int): Maximum results per search type.

    Returns:
        dict: {
            'success': bool,
            'issues': [...],
            'commits': [...],
            'error': str (if failed)
        }
    """
    log_file = Path(__file__).parent / 'host.log'

    if not github_token:
        return {'success': False, 'error': 'GitHub token not configured'}

    if not query or not query.strip():
        return {'success': False, 'error': 'Search query is empty'}

    results = {
        'success': True,
        'issues': [],
        'commits': [],
        'query': query
    }

    headers = {
        'Authorization': f'Bearer {github_token}',
        'Accept': 'application/vnd.github+json',
        'X-GitHub-Api-Version': '2022-11-28'
    }

    try:
        # Build query with optional repo filter
        base_query = query.strip()
        if repos:
            # Parse comma-separated repos and add to query
            repo_list = [r.strip() for r in repos.split(',') if r.strip()]
            if repo_list:
                repo_filter = ' '.join([f'repo:{r}' for r in repo_list])
                base_query = f'{base_query} {repo_filter}'

        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: GitHub search query: {base_query}\n")

        # Search issues
        try:
            issues_url = f'https://api.github.com/search/issues?q={urllib.parse.quote(base_query)}&per_page={max_results}&sort=updated'
            req = urllib.request.Request(issues_url, headers=headers)

            with urllib.request.urlopen(req, timeout=30) as response:
                issues_data = json.loads(response.read().decode('utf-8'))

            for item in issues_data.get('items', [])[:max_results]:
                results['issues'].append({
                    'number': item.get('number'),
                    'title': item.get('title'),
                    'state': item.get('state'),
                    'url': item.get('html_url'),
                    'repo': item.get('repository_url', '').split('/')[-1] if item.get('repository_url') else '',
                    'labels': [l.get('name') for l in item.get('labels', [])],
                    'created_at': item.get('created_at'),
                    'updated_at': item.get('updated_at'),
                    'body_preview': (item.get('body') or '')[:200]
                })

            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: GitHub found {len(results['issues'])} issues\n")

        except urllib.error.HTTPError as e:
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: GitHub issues search error: {e.code} {e.reason}\n")

        # Search commits
        try:
            commits_url = f'https://api.github.com/search/commits?q={urllib.parse.quote(base_query)}&per_page={max_results}&sort=committer-date'
            req = urllib.request.Request(commits_url, headers=headers)

            with urllib.request.urlopen(req, timeout=30) as response:
                commits_data = json.loads(response.read().decode('utf-8'))

            for item in commits_data.get('items', [])[:max_results]:
                commit = item.get('commit', {})
                results['commits'].append({
                    'sha': item.get('sha', '')[:7],
                    'message': commit.get('message', '').split('\n')[0][:100],  # First line only
                    'url': item.get('html_url'),
                    'repo': item.get('repository', {}).get('full_name', ''),
                    'author': commit.get('author', {}).get('name', ''),
                    'date': commit.get('committer', {}).get('date', '')
                })

            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: GitHub found {len(results['commits'])} commits\n")

        except urllib.error.HTTPError as e:
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: GitHub commits search error: {e.code} {e.reason}\n")

        return results

    except urllib.error.HTTPError as e:
        error_msg = f'GitHub API error: {e.code} {e.reason}'
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: {error_msg}\n")
        return {'success': False, 'error': error_msg}
    except Exception as e:
        error_msg = f'GitHub search error: {str(e)}'
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: {error_msg}\n")
        return {'success': False, 'error': error_msg}


def format_github_results_for_context(github_results):
    """
    Format GitHub search results as context text for AI prompts.

    Args:
        github_results (dict): Results from search_github().

    Returns:
        str: Formatted text summarising the GitHub findings.
    """
    if not github_results.get('success'):
        return ''

    parts = []

    issues = github_results.get('issues', [])
    if issues:
        parts.append(f"**Related GitHub Issues ({len(issues)}):**")
        for issue in issues[:5]:  # Limit to 5 for context
            status = '✓' if issue['state'] == 'closed' else '○'
            labels = f" [{', '.join(issue['labels'][:3])}]" if issue['labels'] else ''
            parts.append(f"  {status} #{issue['number']}: {issue['title']}{labels}")
            if issue['body_preview']:
                parts.append(f"    {issue['body_preview'][:100]}...")

    commits = github_results.get('commits', [])
    if commits:
        parts.append(f"\n**Related GitHub Commits ({len(commits)}):**")
        for commit in commits[:5]:  # Limit to 5 for context
            parts.append(f"  • {commit['sha']}: {commit['message']} ({commit['author']})")

    return '\n'.join(parts) if parts else ''


# =============================================================================
# AI Classification with OpenAI
# =============================================================================

def load_topics(project_folder):
    """Load existing topics from topics.json."""
    try:
        topics_file = Path(project_folder) / 'topics.json'
        if topics_file.exists():
            with open(topics_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return data.get('topics', [])
        return []
    except Exception:
        return []


def save_topics(project_folder, topics):
    """Save topics to topics.json."""
    try:
        topics_file = Path(project_folder) / 'topics.json'
        # Deduplicate and sort
        unique_topics = sorted(list(set(topics)))
        with open(topics_file, 'w', encoding='utf-8') as f:
            json.dump({'topics': unique_topics}, f, indent=2)
    except Exception:
        pass


def load_entities(project_folder):
    """Load existing entities (people & roles) from entities.json."""
    try:
        entities_file = Path(project_folder) / 'entities.json'
        if entities_file.exists():
            with open(entities_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return data.get('entities', [])
        return []
    except Exception:
        return []


def save_entities(project_folder, entities):
    """Save entities to entities.json."""
    try:
        entities_file = Path(project_folder) / 'entities.json'
        # Deduplicate and sort
        unique_entities = sorted(list(set(entities)))
        with open(entities_file, 'w', encoding='utf-8') as f:
            json.dump({'entities': unique_entities}, f, indent=2)
    except Exception:
        pass



def classify_entry(entry, api_key, existing_topics, existing_entities, ai_summary=None, custom_prompt=None):
    """
    Classify entry using OpenAI API.
    Returns dict with entity, topics, and people fields.
    ai_summary: Optional AI-generated summary to provide additional context for classification.
    custom_prompt: Optional custom prompt template with placeholders.
    """
    log_file = Path(__file__).parent / 'host.log'

    if not api_key:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: AI classification skipped - no API key\n")
        return None

    try:
        # Build content for classification
        title = entry.get('title', '')
        url = entry.get('url', '')
        notes = entry.get('notes', '')
        selected_text = entry.get('selectedText', '')

        # Combine all available text content
        content_parts = [selected_text, notes]
        if ai_summary:
            content_parts.append(f"AI Description: {ai_summary}")
        content = ' '.join(part for part in content_parts if part).strip()

        # Use prompt from settings
        prompt_template = custom_prompt

        # Replace placeholders in prompt
        prompt = prompt_template.format(
            title=title,
            url=url,
            content=content,
            notes=notes,
            existing_topics=existing_topics,
            existing_people=existing_entities
        )

        # Call OpenAI Responses API with Structured Outputs
        request_data = json.dumps({
            'model': 'gpt-5-nano',
            'input': prompt,
            'text': {
                'format': {
                    'type': 'json_schema',
                    'name': 'classification',
                    'strict': True,
                    'schema': {
                        'type': 'object',
                        'properties': {
                            'entity': {
                                'type': 'string',
                                'enum': ['project', 'task', 'knowledge', 'unclassified']
                            },
                            'topics': {
                                'type': 'array',
                                'items': {'type': 'string'}
                            },
                            'people': {
                                'type': 'array',
                                'items': {'type': 'string'}
                            },
                            'category': {
                                'type': 'string',
                                'enum': ['work', 'personal']
                            },
                            'corrected_notes': {
                                'type': 'string',
                                'description': 'Grammar-corrected version of user notes'
                            }
                        },
                        'required': ['entity', 'topics', 'people', 'category', 'corrected_notes'],
                        'additionalProperties': False
                    }
                }
            }
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=30) as response:
            result = json.loads(response.read().decode('utf-8'))

        # Extract text from response - guaranteed valid JSON with structured outputs
        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        text_content = content_item.get('text', '')
                        classification = json.loads(text_content)
                        with open(log_file, 'a', encoding='utf-8') as f:
                            f.write(f"{datetime.now()}: AI classification result: {classification}\n")
                        return classification

        return None

    except urllib.error.HTTPError as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: AI classification HTTP error: {e.code} {e.reason}\n")
        return None
    except Exception as e:
        import traceback
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: AI classification error: {str(e)}\n")
            f.write(f"{datetime.now()}: AI classification traceback: {traceback.format_exc()}\n")
        return None


# =============================================================================
# AI Summary Functions
# =============================================================================

def summarize_image(image_path, api_key, notes='', custom_prompt=None):
    """Generate summary of an image using GPT-5 vision."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        # Read and encode image
        with open(image_path, 'rb') as f:
            image_data = f.read()
        base64_image = base64.b64encode(image_data).decode('utf-8')

        # Determine MIME type
        ext = Path(image_path).suffix.lower()
        mime_type = 'image/png' if ext == '.png' else 'image/jpeg'

        # Use prompt from settings, format with notes
        prompt = custom_prompt.format(notes=notes if notes else 'None')

        # Call GPT-5 vision API with web search
        request_data = json.dumps({
            'model': 'gpt-5',
            'tools': [{'type': 'web_search'}],
            'input': [{
                'role': 'user',
                'content': [
                    {'type': 'input_text', 'text': prompt},
                    {'type': 'input_image', 'image_url': f'data:{mime_type};base64,{base64_image}'}
                ]
            }]
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=120) as response:
            result = json.loads(response.read().decode('utf-8'))

        # Extract text from response
        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        summary = content_item.get('text', '').strip()
                        with open(log_file, 'a', encoding='utf-8') as f:
                            f.write(f"{datetime.now()}: Image summary generated: {summary[:100]}...\n")
                        return summary

        return None

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: summarize_image error: {str(e)}\n")
        return None


def summarize_with_research(notes, api_key, custom_prompt=None):
    """Generate summary with web research using GPT-5 web search."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        # Use prompt from settings
        prompt = custom_prompt.format(notes=notes)

        request_data = json.dumps({
            'model': 'gpt-5',
            'tools': [{'type': 'web_search'}],
            'input': prompt
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=120) as response:  # Longer timeout for web search
            result = json.loads(response.read().decode('utf-8'))

        # Extract text from response
        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        summary = content_item.get('text', '').strip()
                        with open(log_file, 'a', encoding='utf-8') as f:
                            f.write(f"{datetime.now()}: Research summary generated: {summary[:100]}...\n")
                        return summary

        return None

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: summarize_with_research error: {str(e)}\n")
        return None


def summarize_audio(audio_path, api_key, notes='', custom_prompt=None):
    """Transcribe audio with Whisper, then summarize with speaker identification."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        # Read audio file
        with open(audio_path, 'rb') as f:
            audio_data = f.read()

        # Get filename for API
        filename = Path(audio_path).name

        # Create multipart form data for transcription
        boundary = '----WebKitFormBoundary7MA4YWxkTrZu0gW'
        body = (
            f'--{boundary}\r\n'
            f'Content-Disposition: form-data; name="file"; filename="{filename}"\r\n'
            f'Content-Type: audio/mpeg\r\n\r\n'
        ).encode('utf-8') + audio_data + (
            f'\r\n--{boundary}\r\n'
            f'Content-Disposition: form-data; name="model"\r\n\r\n'
            f'gpt-4o-transcribe\r\n'
            f'--{boundary}--\r\n'
        ).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/audio/transcriptions',
            data=body,
            headers={
                'Content-Type': f'multipart/form-data; boundary={boundary}',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=120) as response:
            transcription_result = json.loads(response.read().decode('utf-8'))

        transcript = transcription_result.get('text', '')

        if not transcript:
            return None

        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: Audio transcribed: {transcript[:100]}...\n")

        # Build prompt using template from settings
        prompt = custom_prompt.format(notes=notes if notes else 'None', transcript=transcript)

        # Use structured outputs for guaranteed JSON response
        summary_request = json.dumps({
            'model': 'gpt-5-nano',
            'input': prompt,
            'text': {
                'format': {
                    'type': 'json_schema',
                    'name': 'audio_analysis',
                    'strict': True,
                    'schema': {
                        'type': 'object',
                        'properties': {
                            'summary': {
                                'type': 'string',
                                'description': 'Comprehensive summary of what is discussed/happening'
                            },
                            'speakers': {
                                'type': 'array',
                                'items': {'type': 'string'},
                                'description': 'List of speakers identified (Speaker 1, Speaker 2, etc)'
                            },
                            'transcript': {
                                'type': 'string',
                                'description': 'The full transcript'
                            }
                        },
                        'required': ['summary', 'speakers', 'transcript'],
                        'additionalProperties': False
                    }
                }
            }
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=summary_request,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=30) as response:
            result = json.loads(response.read().decode('utf-8'))

        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        # Return JSON string (structured output)
                        summary = content_item.get('text', '').strip()
                        with open(log_file, 'a', encoding='utf-8') as f:
                            f.write(f"{datetime.now()}: Audio summary generated (structured): {summary[:100]}...\n")
                        return summary

        return None

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: summarize_audio error: {str(e)}\n")
        return None


def summarize_document(file_path, entry_type, api_key, custom_prompt=None):
    """Summarize a document (PDF, markdown, Office files)."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        content = None

        # Use prompt from settings
        prompt_template = custom_prompt

        # Extract text based on document type
        if entry_type == 'markdown':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Markdown text read: {len(content)} chars\n")

        elif entry_type == 'pdf':
            try:
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text_parts = []
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                    content = '\n'.join(text_parts)
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: PDF text extracted with PyPDF2: {len(content)} chars\n")
            except ImportError:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: PyPDF2 not installed - run: py -m pip install PyPDF2\n")
                return None
            except Exception as e:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: PyPDF2 extraction failed: {str(e)}\n")
                return None

        elif entry_type in ['ms-word', 'docx']:
            try:
                from docx import Document
                doc = Document(file_path)
                content = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Word text extracted with python-docx: {len(content)} chars\n")
            except ImportError:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: python-docx not installed - run: py -m pip install python-docx\n")
                return None
            except Exception as e:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Word extraction failed: {str(e)}\n")
                return None

        elif entry_type in ['ms-excel', 'xlsx']:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, data_only=True)
                text_parts = []
                for sheet in wb.worksheets:
                    text_parts.append(f"Sheet: {sheet.title}")
                    for row in sheet.iter_rows(max_row=100, values_only=True):  # Limit rows
                        row_text = ' | '.join([str(cell) for cell in row if cell is not None])
                        if row_text.strip():
                            text_parts.append(row_text)
                content = '\n'.join(text_parts)
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Excel text extracted with openpyxl: {len(content)} chars\n")
            except ImportError:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: openpyxl not installed - run: py -m pip install openpyxl\n")
                return None
            except Exception as e:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Excel extraction failed: {str(e)}\n")
                return None

        elif entry_type in ['ms-powerpoint', 'pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_text.append(shape.text)
                    if slide_text:
                        text_parts.append(f"Slide {slide_num}: " + ' | '.join(slide_text))
                content = '\n'.join(text_parts)
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: PowerPoint text extracted with python-pptx: {len(content)} chars\n")
            except ImportError:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: python-pptx not installed - run: py -m pip install python-pptx\n")
                return None
            except Exception as e:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: PowerPoint extraction failed: {str(e)}\n")
                return None

        else:
            # Fallback: try to read as plain text
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: {entry_type} read as plain text: {len(content)} chars\n")
            except Exception as e:
                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Could not read {entry_type} file: {str(e)}\n")
                return None

        if not content or len(content.strip()) < 10:
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Document text extraction yielded no content\n")
            return None

        # Summarize using gpt-5-nano with structured outputs
        prompt = prompt_template.format(content=content[:8000])
        request_data = json.dumps({
            'model': 'gpt-5-nano',
            'input': prompt,
            'text': {
                'format': {
                    'type': 'json_schema',
                    'name': 'document_summary',
                    'strict': True,
                    'schema': {
                        'type': 'object',
                        'properties': {
                            'brief': {
                                'type': 'string',
                                'description': '1-2 sentence summary'
                            },
                            'paragraph': {
                                'type': 'string',
                                'description': 'Single paragraph summary'
                            },
                            'detailed': {
                                'type': 'string',
                                'description': 'Detailed multi-paragraph summary'
                            },
                            'entities': {
                                'type': 'object',
                                'properties': {
                                    'people': {
                                        'type': 'array',
                                        'items': {'type': 'string'},
                                        'description': 'Key people mentioned'
                                    },
                                    'topics': {
                                        'type': 'array',
                                        'items': {'type': 'string'},
                                        'description': 'Key topics'
                                    },
                                    'places': {
                                        'type': 'array',
                                        'items': {'type': 'string'},
                                        'description': 'Key places'
                                    },
                                    'other': {
                                        'type': 'array',
                                        'items': {'type': 'string'},
                                        'description': 'Other relevant entities'
                                    }
                                },
                                'required': ['people', 'topics', 'places', 'other'],
                                'additionalProperties': False
                            }
                        },
                        'required': ['brief', 'paragraph', 'detailed', 'entities'],
                        'additionalProperties': False
                    }
                }
            }
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=30) as response:
            result = json.loads(response.read().decode('utf-8'))

        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        # Return JSON string (structured output)
                        summary = content_item.get('text', '').strip()
                        with open(log_file, 'a', encoding='utf-8') as f:
                            f.write(f"{datetime.now()}: Document summary generated (structured): {summary[:100]}...\n")
                        return summary

        return None

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: summarize_document error: {str(e)}\n")
        return None


def summarize_link(entry, api_key, custom_prompt=None, text_prompt=None):
    """
    Summarize a web link by browsing and analysing its content.

    Uses GPT-5 with web search capability to fetch and understand the
    linked page, then generates a comprehensive summary with sources.

    Args:
        entry (dict): Entry containing 'url', 'title', and optional 'notes'.
        api_key (str): OpenAI API key.
        custom_prompt (str|None): Custom prompt template for link summary.
        text_prompt (str|None): Custom prompt for text fallback.

    Returns:
        str|None: Summary text with sources appended, or None on failure.

    API configuration:
        - Model: gpt-5
        - Tools: web_search (to fetch page content)
        - Reasoning effort: medium
        - Timeout: 90 seconds
    """
    log_file = Path(__file__).parent / 'host.log'

    try:
        title = entry.get('title', '')
        url = entry.get('url', '')
        notes = entry.get('notes', '')

        if not url:
            # No URL to browse - fall back to basic summary
            return summarize_text(entry, api_key, text_prompt)

        # Use prompt from settings
        prompt = custom_prompt.format(url=url, title=title, notes=notes if notes else 'None')

        # Use structured outputs with web search
        request_data = json.dumps({
            'model': 'gpt-5',
            'tools': [{'type': 'web_search'}],
            'include': ['web_search_call.action.sources'],
            'input': prompt,
            'text': {
                'format': {
                    'type': 'json_schema',
                    'name': 'link_summary',
                    'strict': True,
                    'schema': {
                        'type': 'object',
                        'properties': {
                            'brief': {
                                'type': 'string',
                                'description': '2-3 sentence summary'
                            },
                            'detailed': {
                                'type': 'string',
                                'description': 'Detailed multi-paragraph analysis'
                            },
                            'key_points': {
                                'type': 'array',
                                'items': {'type': 'string'},
                                'description': 'Key information, facts, takeaways'
                            },
                            'context': {
                                'type': 'string',
                                'description': 'Relevant context, related links, supporting evidence'
                            },
                            'sources': {
                                'type': 'array',
                                'items': {'type': 'string'},
                                'description': 'Source URLs'
                            }
                        },
                        'required': ['brief', 'detailed', 'key_points', 'context', 'sources'],
                        'additionalProperties': False
                    }
                }
            }
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=240) as response:  # 4 min timeout for web search + reasoning
            result = json.loads(response.read().decode('utf-8'))

        # Extract structured JSON response
        summary_text = None

        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        summary_text = content_item.get('text', '').strip()

        if summary_text:
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Link summary generated (web+reasoning): {summary_text[:100]}...\n")
            return summary_text

        return None

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: summarize_link error: {str(e)}\n")
        return None


def summarize_text(entry, api_key, custom_prompt=None):
    """Summarize text content (snippets, notes, ideas, paragraphs) with web search."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        notes = entry.get('notes', '')
        selected_text = entry.get('selectedText', '')

        # Need at least one of selectedText or notes
        if not selected_text and not notes:
            return None

        # Use prompt from settings - send BOTH selectedText and notes
        prompt = custom_prompt.format(
            selectedText=selected_text[:2000] if selected_text else 'None',
            notes=notes[:500] if notes else 'None'
        )

        # Use gpt-5 with web search for richer summaries
        request_data = json.dumps({
            'model': 'gpt-5',
            'tools': [{'type': 'web_search'}],
            'input': prompt
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=90) as response:  # Longer timeout for web search
            result = json.loads(response.read().decode('utf-8'))

        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        summary = content_item.get('text', '').strip()
                        with open(log_file, 'a', encoding='utf-8') as f:
                            f.write(f"{datetime.now()}: Text summary generated: {summary[:100]}...\n")
                        return summary

        return None

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: summarize_text error: {str(e)}\n")
        return None


def generate_ai_summary(entry_type, entry, file_path, api_key, prompts=None):
    """
    Route to the appropriate AI summary function based on entry type.

    Dispatches to specialised summary functions depending on content type:
    - Images/screenshots: Vision analysis
    - Audio: Whisper transcription + analysis
    - Documents: Text extraction + summarisation
    - Links: Web browsing + summarisation
    - Text: Direct summarisation

    Args:
        entry_type (str): Type of entry ('screenshot', 'audio', 'pdf', etc.).
        entry (dict): Entry data with notes, selectedText, url, etc.
        file_path (str|None): Path to associated file, if any.
        api_key (str): OpenAI API key.
        prompts (dict|None): Custom prompts for each summary type.

    Returns:
        str|None: Generated summary text, or None if skipped/failed.
    """
    log_file = Path(__file__).parent / 'host.log'
    prompts = prompts or {}

    try:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: Generating AI summary for type: {entry_type}\n")

        # Skip types that don't need summary
        if entry_type in NO_SUMMARY_TYPES:
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Skipping AI summary for type: {entry_type}\n")
            return None

        if entry_type in IMAGE_TYPES:
            if file_path:
                notes = entry.get('notes', '')
                return summarize_image(file_path, api_key, notes, prompts.get('image'))
        elif entry_type in RESEARCH_TYPES:
            notes = entry.get('notes', '')
            if notes:
                return summarize_with_research(notes, api_key, prompts.get('research'))
        elif entry_type in AUDIO_TYPES:
            if file_path:
                notes = entry.get('notes', '')
                return summarize_audio(file_path, api_key, notes, prompts.get('audio'))
        elif entry_type in DOCUMENT_TYPES:
            if file_path:
                return summarize_document(file_path, entry_type, api_key, prompts.get('document'))
        elif entry_type in TEXT_TYPES:
            return summarize_text(entry, api_key, prompts.get('text'))
        elif entry_type in LINK_TYPES:
            return summarize_link(entry, api_key, prompts.get('link'), prompts.get('text'))
        else:
            # Default: try text summary for unknown types
            return summarize_text(entry, api_key, prompts.get('text'))

        return None

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: generate_ai_summary error: {str(e)}\n")
        return None


def add_summary_to_entry(project_folder, timestamp, summary):
    """Add AI Summary metadata line to an existing kb.md entry."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return False

        folder_path = result
        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return False

        with open(kb_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # Find entry with matching timestamp
        i = 0
        while i < len(lines):
            if f'`{timestamp}`' in lines[i]:
                # Found entry - find where to insert summary (after classification)
                j = i + 1
                insert_pos = j

                while j < len(lines):
                    if lines[j].startswith('- '):
                        # Hit next entry
                        insert_pos = j
                        break
                    elif lines[j].strip() == '':
                        # Hit blank line (end of entry)
                        insert_pos = j
                        break
                    j += 1
                else:
                    insert_pos = len(lines)

                # Clean up summary (remove newlines, keep full content)
                clean_summary = ' '.join(summary.split())

                # Insert summary line
                lines.insert(insert_pos, f"  - AI Summary: {clean_summary}\n")

                # Write back
                with open(kb_file, 'w', encoding='utf-8') as f:
                    f.writelines(lines)

                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: AI Summary added for {timestamp}\n")

                return True

            i += 1

        return False

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: add_summary_to_entry error: {str(e)}\n")
        return False


def read_message():
    """Read a message from stdin (sent by extension)."""
    # Read the message length (first 4 bytes)
    raw_length = sys.stdin.buffer.read(4)
    if not raw_length:
        return None

    # Unpack message length
    message_length = struct.unpack('=I', raw_length)[0]

    # Read the message
    message = sys.stdin.buffer.read(message_length).decode('utf-8')

    return json.loads(message)


def send_message(message):
    """Send a message to stdout (back to extension)."""
    encoded_message = json.dumps(message).encode('utf-8')
    encoded_length = struct.pack('=I', len(encoded_message))

    sys.stdout.buffer.write(encoded_length)
    sys.stdout.buffer.write(encoded_message)
    sys.stdout.buffer.flush()


def format_markdown_entry(entry, screenshot_path=None, file_path=None, classification=None):
    """Format entry data as markdown with consistent structure.

    Format: - `type` | `source` | `timestamp` | [Title](URL) | `group: Name (color)`
    """
    lines = []

    # Extract fields
    title = entry.get('title', 'Untitled')
    url = entry.get('url', '')  # URL is now separate from source
    entry_type = entry['type']
    timestamp = entry['captured']
    source = entry.get('source', 'browser')  # 'browser' or 'widget'

    # Build main line: type | source | timestamp | title (with optional URL) | optional group
    main_line = f"- `{entry_type}` | `{source}` | `{timestamp}` | "

    # Add title with URL only if URL exists and is a valid URL (not a filename)
    if url and url.startswith(('http://', 'https://', 'file://', 'chrome-extension://')):
        main_line += f"[{title}]({url})"
    else:
        main_line += title

    # Add tab group if present
    tab_group = entry.get('tabGroup')
    if tab_group:
        group_name = tab_group.get('groupName', '')
        group_color = tab_group.get('groupColor', '')

        if group_name:
            main_line += f" | `group: {group_name} ({group_color})`"
        elif group_color:
            main_line += f" | `group: ({group_color})`"

    lines.append(main_line)

    # Add selected text as blockquote (if present)
    selected_text = entry.get('selectedText', '').strip()
    if selected_text:
        # Clean up and format as blockquote
        text_lines = [line.strip() for line in selected_text.split('\n') if line.strip()]
        for line in text_lines:
            lines.append(f"  - > {line}")

    # Add screenshot image reference
    if screenshot_path:
        lines.append(f"  - ![Screenshot]({screenshot_path})")

    # Add file attachment link
    if file_path:
        lines.append(f"  - [Attachment]({file_path})")

    # Add notes (user commentary) if present
    notes = entry.get('notes', '').strip()
    if notes:
        # Clean up and prefix with "Notes:"
        note_lines = [line.strip() for line in notes.split('\n') if line.strip()]
        if len(note_lines) == 1:
            lines.append(f"  - Notes: {note_lines[0]}")
        else:
            lines.append(f"  - Notes: {note_lines[0]}")
            for line in note_lines[1:]:
                lines.append(f"    {line}")

    # Add relationship fields if present
    parent_id = entry.get('parentId')
    if parent_id:
        lines.append(f"  - ParentId: {parent_id}")

    related = entry.get('related')
    if related:
        # Format: timestamp (type), timestamp (type)
        lines.append(f"  - Related: {related}")

    similar = entry.get('similar')
    if similar:
        # Format: timestamp (score), timestamp (score)
        lines.append(f"  - Similar: {similar}")

    # Add page metadata if present (optional fields)
    page_meta = entry.get('pageMetadata')
    if page_meta:
        desc = page_meta.get('description', '').strip()
        og_image = page_meta.get('ogImage', '').strip()
        author = page_meta.get('author', '').strip()
        pub_date = page_meta.get('publishedDate', '').strip()
        read_time = page_meta.get('readingTime')

        if desc:
            # Truncate long descriptions
            if len(desc) > 200:
                desc = desc[:197] + '...'
            lines.append(f"  - Description: {desc}")
        if og_image:
            lines.append(f"  - Image: {og_image}")
        if author:
            lines.append(f"  - Author: {author}")
        if pub_date:
            lines.append(f"  - Published: {pub_date}")
        if read_time and read_time > 0:
            lines.append(f"  - ReadTime: {read_time} min")

    # Add AI classification metadata if present
    if classification:
        entity = classification.get('entity', '')
        topics = classification.get('topics', [])
        people = classification.get('people', [])

        if entity:
            lines.append(f"  - Entity: {entity}")
        if topics:
            lines.append(f"  - Topics: {', '.join(topics)}")
        if people:
            lines.append(f"  - People: {', '.join(people)}")

    lines.append('')  # Single blank line between entries

    return '\n'.join(lines)


def save_screenshot(project_folder, screenshot_data_url, timestamp):
    """Save screenshot image to file."""
    log_file = Path(__file__).parent / 'host.log'
    try:
        # Validate project folder
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Screenshot validation error: {result}\n")
            return None

        folder_path = result  # result is the validated Path object
        screenshots_folder = folder_path / 'screenshots'
        screenshots_folder.mkdir(exist_ok=True)

        # Extract base64 data from data URL (format: data:image/png;base64,...)
        if ',' in screenshot_data_url:
            base64_data = screenshot_data_url.split(',', 1)[1]
        else:
            base64_data = screenshot_data_url

        # Decode base64 to binary
        image_data = base64.b64decode(base64_data)

        # Generate filename from timestamp
        filename = f"screenshot_{timestamp.replace(' ', '_').replace(':', '-')}.png"
        screenshot_path = screenshots_folder / filename

        # Save image file
        with open(screenshot_path, 'wb') as f:
            f.write(image_data)

        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: Screenshot saved: {screenshot_path}\n")

        # Return relative path for markdown
        return f"screenshots/{filename}"

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: Screenshot error: {str(e)}\n")
        return None


def save_file(project_folder, file_data_url, filename, timestamp):
    """Save uploaded file to /files folder."""
    log_file = Path(__file__).parent / 'host.log'
    try:
        # Validate project folder
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: File save validation error: {result}\n")
            return None

        folder_path = result  # result is the validated Path object
        files_folder = folder_path / 'files'
        files_folder.mkdir(exist_ok=True)

        # Extract base64 data from data URL (format: data:mime/type;base64,...)
        if ',' in file_data_url:
            base64_data = file_data_url.split(',', 1)[1]
        else:
            base64_data = file_data_url

        # Decode base64 to binary
        file_data = base64.b64decode(base64_data)

        # Sanitize filename to prevent path traversal
        safe_base = sanitize_filename(filename)

        # Generate filename with timestamp to avoid collisions
        name_parts = safe_base.rsplit('.', 1)
        if len(name_parts) == 2:
            base_name, extension = name_parts
            safe_filename = f"{base_name}_{timestamp.replace(' ', '_').replace(':', '-')}.{extension}"
        else:
            safe_filename = f"{safe_base}_{timestamp.replace(' ', '_').replace(':', '-')}"

        file_path = files_folder / safe_filename

        # Save file
        with open(file_path, 'wb') as f:
            f.write(file_data)

        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: File saved: {file_path}\n")

        # Return relative path for markdown
        return f"files/{safe_filename}"

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: File save error: {str(e)}\n")
        return None


def append_to_kb(project_folder, entry, api_key=None):
    """Append entry to kb.md file (at the top) with optional AI classification."""
    try:
        # Validate project folder
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return {'success': False, 'error': result}

        folder_path = result  # result is the validated Path object

        # Validate entry structure
        is_valid, error = validate_entry(entry)
        if not is_valid:
            return {'success': False, 'error': error}

        kb_file = folder_path / 'kb.md'

        # Handle screenshot if present (browser screenshot capture)
        screenshot_path = None
        if entry.get('type') == 'screenshot' and 'screenshot' in entry:
            screenshot_path = save_screenshot(
                project_folder,
                entry['screenshot'],
                entry['captured']
            )
            if not screenshot_path:
                return {'success': False, 'error': 'Failed to save screenshot'}

        # Handle file data if present (files, audio, images, etc.)
        file_path = None
        if 'fileData' in entry:
            file_path = save_file(
                project_folder,
                entry['fileData'],
                entry['title'],  # Use title as filename
                entry['captured']
            )
            if not file_path:
                return {'success': False, 'error': 'Failed to save file'}

        # Format entry WITHOUT classification (classification runs in background via separate call)
        new_content = format_markdown_entry(entry, screenshot_path, file_path, classification=None)

        # Read existing content if file exists
        existing_content = ''
        if kb_file.exists():
            with open(kb_file, 'r', encoding='utf-8') as f:
                existing_content = f.read()

        # Write new entry at top, followed by existing content
        with open(kb_file, 'w', encoding='utf-8') as f:
            f.write(new_content)
            if existing_content:
                f.write(existing_content)

        # Return success and background task info
        # main() will send response first, then run background processing
        # Include file path for AI summary (screenshot or uploaded file)
        saved_file_path = screenshot_path or file_path  # Relative path like "screenshots/file.png"
        return {
            'success': True,
            'file': str(kb_file),
            '_background_task': {
                'project_folder': project_folder,
                'timestamp': entry['captured'],
                'entry': entry,
                'api_key': api_key,
                'file_path': str(folder_path / saved_file_path) if saved_file_path else None
            } if api_key else None
        }

    except Exception as e:
        return {'success': False, 'error': str(e)}


def update_last_entry(project_folder, timestamp, new_content):
    """Update the content of the most recent entry with matching timestamp."""
    try:
        # Validate project folder
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return {'success': False, 'error': result}

        folder_path = result  # result is the validated Path object

        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return {'success': False, 'error': 'kb.md file does not exist'}

        # Read the entire file
        with open(kb_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # Find the entry with matching timestamp and update its content
        updated = False
        i = 0
        while i < len(lines):
            line = lines[i]

            # Check if this line contains the timestamp we're looking for
            if f'`{timestamp}`' in line:
                # Found the entry - now look for content lines (indented with '  - ')
                # Skip until we find content or hit the next entry
                j = i + 1
                content_found = False

                # Look ahead to find where content starts
                while j < len(lines):
                    if lines[j].startswith('  - ') and not lines[j].startswith('  - ![') and not lines[j].startswith('  - [📎'):
                        # Found content line - replace it
                        if not content_found:
                            lines[j] = f'  - {new_content}\n'
                            content_found = True
                            updated = True
                        else:
                            # Remove additional old content lines
                            lines.pop(j)
                            continue
                    elif lines[j].startswith('- ') or (j + 1 < len(lines) and lines[j].strip() == ''):
                        # Hit next entry or blank line separator
                        break
                    j += 1

                # If no content was found, insert it after the main line
                if not content_found and new_content.strip():
                    # Find where to insert (after screenshot/file link if present)
                    insert_pos = i + 1
                    while insert_pos < len(lines) and (lines[insert_pos].startswith('  - ![') or lines[insert_pos].startswith('  - [📎')):
                        insert_pos += 1
                    lines.insert(insert_pos, f'  - {new_content}\n')
                    updated = True

                break

            i += 1

        if updated:
            # Write back to file
            with open(kb_file, 'w', encoding='utf-8') as f:
                f.writelines(lines)
            return {'success': True, 'message': 'Entry updated'}
        else:
            return {'success': False, 'error': f'Entry with timestamp {timestamp} not found'}

    except Exception as e:
        return {'success': False, 'error': str(e)}


def update_entry_notes(project_folder, timestamp, new_notes):
    """Update the Notes: line of an entry with matching timestamp."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        # Validate project folder
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return False

        folder_path = result
        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return False

        with open(kb_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # Find entry with matching timestamp
        updated = False
        i = 0
        while i < len(lines):
            if f'`{timestamp}`' in lines[i]:
                # Found entry - look for Notes: line and any continuation lines
                j = i + 1
                notes_start = -1
                notes_end = -1

                while j < len(lines):
                    if lines[j].startswith('  - Notes:'):
                        notes_start = j
                        notes_end = j + 1
                        # Find continuation lines (4-space indent, not starting with "  - ")
                        while notes_end < len(lines):
                            if lines[notes_end].startswith('    ') and not lines[notes_end].startswith('  - '):
                                notes_end += 1
                            else:
                                break
                        break
                    elif lines[j].startswith('- ') or lines[j].strip() == '':
                        # Hit next entry or blank line - no notes found
                        break
                    j += 1

                # Build new notes lines
                note_lines_list = [line.strip() for line in new_notes.split('\n') if line.strip()]
                new_note_lines = []
                if note_lines_list:
                    new_note_lines.append(f"  - Notes: {note_lines_list[0]}\n")
                    for line in note_lines_list[1:]:
                        new_note_lines.append(f"    {line}\n")

                if notes_start >= 0:
                    # Replace existing notes (including continuation lines)
                    lines[notes_start:notes_end] = new_note_lines
                    updated = True
                elif new_notes.strip():
                    # Insert new notes after selected text, screenshots, attachments
                    insert_pos = i + 1
                    while insert_pos < len(lines):
                        line = lines[insert_pos]
                        if line.startswith('  - >') or line.startswith('  - ![') or line.startswith('  - ['):
                            insert_pos += 1
                        else:
                            break
                    for idx, note_line in enumerate(new_note_lines):
                        lines.insert(insert_pos + idx, note_line)
                    updated = True

                break
            i += 1

        if updated:
            with open(kb_file, 'w', encoding='utf-8') as f:
                f.writelines(lines)
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Notes updated for {timestamp}\n")
            return True

        return False

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: update_entry_notes error: {str(e)}\n")
        return False


def truncate_entry_for_ai(entry, max_notes=200, max_summary=200):
    """Prepare entry for AI comparison prompt with truncated fields."""
    return {
        'timestamp': entry.get('timestamp', ''),
        'title': entry.get('title', '')[:100],
        'entity': entry.get('entity', ''),
        'topics': entry.get('topics', [])[:5],
        'notes': entry.get('content', '')[:max_notes],
        'summary': entry.get('aiSummary', '')[:max_summary]
    }


def load_all_entries(project_folder):
    """Load and parse all entries from kb.md."""
    log_file = Path(__file__).parent / 'host.log'
    try:
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return []

        folder_path = result
        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return []

        with open(kb_file, 'r', encoding='utf-8') as f:
            content = f.read()

        # Simple parsing - extract timestamp, title, entity, topics, notes from each entry
        entries = []
        current_entry = None

        for line in content.split('\n'):
            # New entry line
            if line.startswith('- `') and '|' in line:
                if current_entry:
                    entries.append(current_entry)

                # Parse: - `type` | `source` | `timestamp` | Title
                parts = line.split('|')
                if len(parts) >= 3:
                    timestamp = parts[2].strip().strip('`')
                    title_part = parts[3].strip() if len(parts) > 3 else ''
                    # Extract title from [Title](url) or plain title
                    if title_part.startswith('['):
                        title = title_part.split(']')[0][1:]
                    else:
                        title = title_part

                    current_entry = {
                        'timestamp': timestamp,
                        'title': title,
                        'entity': '',
                        'topics': [],
                        'people': [],
                        'content': '',
                        'aiSummary': ''
                    }
            elif current_entry and line.startswith('  - '):
                content_line = line[4:]
                if content_line.startswith('Notes: '):
                    current_entry['content'] = content_line[7:]
                elif content_line.startswith('Entity: '):
                    current_entry['entity'] = content_line[8:]
                elif content_line.startswith('Topics: '):
                    current_entry['topics'] = [t.strip() for t in content_line[8:].split(',')]
                elif content_line.startswith('People: '):
                    current_entry['people'] = [p.strip() for p in content_line[8:].split(',')]
                elif content_line.startswith('AI Summary: '):
                    current_entry['aiSummary'] = content_line[12:]

        if current_entry:
            entries.append(current_entry)

        return entries

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: load_all_entries error: {str(e)}\n")
        return []


def fuzzy_search_candidates(all_entries, target_description, target_type, temporal_hint, current_timestamp, max_candidates=50):
    """
    Find candidate entries that might match a reference using fuzzy keyword scoring.

    Args:
        all_entries: List of all parsed entries
        target_description: What the user is referring to (e.g., "ultrathink project")
        target_type: Entity type filter (project/task/knowledge/any)
        temporal_hint: Time filter (recent/yesterday/last_week/last_month/none)
        current_timestamp: Current entry timestamp for relative time calculations
        max_candidates: Maximum candidates to return

    Returns:
        List of candidate entries sorted by relevance score
    """
    from datetime import datetime, timedelta

    target_lower = target_description.lower()
    target_words = set(target_lower.split())
    candidates = []

    # Parse current timestamp for temporal filtering
    try:
        current_dt = datetime.strptime(current_timestamp, '%Y-%m-%d %H:%M:%S')
    except:
        current_dt = datetime.now()

    for entry in all_entries:
        # Skip self
        if entry.get('timestamp') == current_timestamp:
            continue

        score = 0

        # Filter by entity type if specified
        if target_type and target_type != 'any':
            if entry.get('entity') != target_type:
                continue

        # Apply temporal filter
        if temporal_hint and temporal_hint != 'none':
            try:
                entry_dt = datetime.strptime(entry['timestamp'], '%Y-%m-%d %H:%M:%S')
                days_ago = (current_dt - entry_dt).days

                if temporal_hint == 'yesterday' and days_ago != 1:
                    continue
                elif temporal_hint == 'recent' and days_ago > 7:
                    continue
                elif temporal_hint == 'last_week' and days_ago > 14:
                    continue
                elif temporal_hint == 'last_month' and days_ago > 45:
                    continue
            except:
                pass

        # Score by title match (highest weight)
        title_lower = entry.get('title', '').lower()
        if target_lower in title_lower:
            score += 50
        else:
            title_words = set(title_lower.split())
            overlap = len(target_words & title_words)
            score += overlap * 15

        # Score by topic match
        for topic in entry.get('topics', []):
            if topic.lower() in target_lower or target_lower in topic.lower():
                score += 25

        # Score by notes content match
        notes_lower = entry.get('content', '').lower()
        if target_lower in notes_lower:
            score += 20
        else:
            notes_words = set(notes_lower.split())
            overlap = len(target_words & notes_words)
            score += overlap * 5

        if score > 0:
            candidates.append({'entry': entry, 'score': score})

    # Sort by score descending
    candidates.sort(key=lambda x: x['score'], reverse=True)
    return [c['entry'] for c in candidates[:max_candidates]]


def extract_relationships(entry, api_key, custom_prompt):
    """
    Step 5a: Use AI to extract relationship references from notes.

    Returns dict with 'references' list and 'has_relationships' boolean.
    """
    log_file = Path(__file__).parent / 'host.log'

    notes = entry.get('notes', '').strip()
    if not notes or len(notes) < 10:
        return {'references': [], 'has_relationships': False}

    try:
        prompt = custom_prompt.format(
            notes=notes,
            title=entry.get('title', ''),
            topics=', '.join(entry.get('topics', []))
        )

        # Call OpenAI Responses API with Structured Outputs
        request_data = json.dumps({
            'model': 'gpt-5-nano',
            'input': prompt,
            'text': {
                'format': {
                    'type': 'json_schema',
                    'name': 'relationship_extraction',
                    'strict': True,
                    'schema': {
                        'type': 'object',
                        'properties': {
                            'references': {
                                'type': 'array',
                                'items': {
                                    'type': 'object',
                                    'properties': {
                                        'phrase': {'type': 'string'},
                                        'target': {'type': 'string'},
                                        'type': {'type': 'string', 'enum': ['project', 'task', 'knowledge', 'any']},
                                        'temporal': {'type': 'string', 'enum': ['recent', 'yesterday', 'last_week', 'last_month', 'none']}
                                    },
                                    'required': ['phrase', 'target', 'type', 'temporal'],
                                    'additionalProperties': False
                                }
                            },
                            'has_relationships': {'type': 'boolean'}
                        },
                        'required': ['references', 'has_relationships'],
                        'additionalProperties': False
                    }
                }
            }
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=30) as response:
            result = json.loads(response.read().decode('utf-8'))

        # Extract text from response - guaranteed valid JSON with structured outputs
        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        text_content = content_item.get('text', '')
                        return json.loads(text_content)

        return {'references': [], 'has_relationships': False}

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: extract_relationships error: {str(e)}\n")
        return {'references': [], 'has_relationships': False}


def resolve_with_ai(phrase, target, candidates, api_key, custom_prompt):
    """
    Step 5b: Use AI to pick the best match from fuzzy search candidates.

    Returns dict with 'timestamp' (or null) and 'confidence'.
    """
    log_file = Path(__file__).parent / 'host.log'

    if not candidates:
        return {'timestamp': None, 'confidence': 'none'}

    try:
        # Format candidates for prompt (truncated)
        candidates_text = ""
        for i, entry in enumerate(candidates[:20], 1):  # Limit to top 20
            truncated = truncate_entry_for_ai(entry)
            candidates_text += f"\n[{i}] {truncated['timestamp']} | {truncated['title']}"
            if truncated['entity']:
                candidates_text += f" | {truncated['entity']}"
            if truncated['topics']:
                candidates_text += f" | Topics: {', '.join(truncated['topics'][:3])}"
            if truncated['notes']:
                candidates_text += f"\n    Notes: {truncated['notes'][:150]}..."

        prompt = custom_prompt.format(
            phrase=phrase,
            target=target,
            candidates=candidates_text
        )

        # Call OpenAI Responses API with Structured Outputs
        request_data = json.dumps({
            'model': 'gpt-5-nano',
            'input': prompt,
            'text': {
                'format': {
                    'type': 'json_schema',
                    'name': 'relationship_resolution',
                    'strict': True,
                    'schema': {
                        'type': 'object',
                        'properties': {
                            'timestamp': {'type': ['string', 'null']},
                            'confidence': {'type': 'string', 'enum': ['high', 'medium', 'low', 'none']}
                        },
                        'required': ['timestamp', 'confidence'],
                        'additionalProperties': False
                    }
                }
            }
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=30) as response:
            result = json.loads(response.read().decode('utf-8'))

        # Extract text from response - guaranteed valid JSON with structured outputs
        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        text_content = content_item.get('text', '')
                        return json.loads(text_content)

        return {'timestamp': None, 'confidence': 'none'}

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: resolve_with_ai error: {str(e)}\n")
        return {'timestamp': None, 'confidence': 'none'}


def get_similarity_candidates(all_entries, current_entry, max_candidates=100):
    """
    Pre-filter entries for content similarity comparison.

    Prioritizes entries with:
    - Same topics (highest)
    - Same entity type
    - Same people
    - Recent entries (last 30 days)
    """
    from datetime import datetime, timedelta

    current_timestamp = current_entry.get('timestamp', '')
    current_topics = set(current_entry.get('topics', []))
    current_entity = current_entry.get('entity', '')
    current_people = set(current_entry.get('people', []))

    try:
        current_dt = datetime.strptime(current_timestamp, '%Y-%m-%d %H:%M:%S')
    except:
        current_dt = datetime.now()

    candidates = []

    for entry in all_entries:
        # Skip self
        if entry.get('timestamp') == current_timestamp:
            continue

        score = 0
        entry_topics = set(entry.get('topics', []))
        entry_people = set(entry.get('people', []))

        # Topic overlap (highest priority)
        topic_overlap = len(current_topics & entry_topics)
        score += topic_overlap * 30

        # Same entity type
        if entry.get('entity') == current_entity and current_entity:
            score += 15

        # People overlap
        people_overlap = len(current_people & entry_people)
        score += people_overlap * 25

        # Recency bonus
        try:
            entry_dt = datetime.strptime(entry['timestamp'], '%Y-%m-%d %H:%M:%S')
            days_ago = (current_dt - entry_dt).days
            if days_ago <= 7:
                score += 20
            elif days_ago <= 30:
                score += 10
        except:
            pass

        if score > 0:
            candidates.append({'entry': entry, 'score': score})

    # Sort by score and return top candidates
    candidates.sort(key=lambda x: x['score'], reverse=True)
    return [c['entry'] for c in candidates[:max_candidates]]


def find_similar_entries(entry, candidates, api_key, custom_prompt):
    """
    Step 6: Use AI to find semantically similar entries.

    Returns list of {'timestamp', 'score', 'reason'} for similar entries.
    """
    log_file = Path(__file__).parent / 'host.log'

    if not candidates:
        return []

    try:
        # Format new entry
        new_entry_text = f"""Title: {entry.get('title', '')}
Topics: {', '.join(entry.get('topics', []))}
Entity: {entry.get('entity', '')}
Notes: {entry.get('notes', '')[:300]}
Summary: {entry.get('aiSummary', '')[:300]}"""

        # Format candidates (truncated)
        candidates_text = ""
        for i, cand in enumerate(candidates[:50], 1):  # Limit to 50
            truncated = truncate_entry_for_ai(cand, max_notes=150, max_summary=150)
            candidates_text += f"\n[{i}] {truncated['timestamp']} | {truncated['title']}"
            if truncated['entity']:
                candidates_text += f" ({truncated['entity']})"
            if truncated['topics']:
                candidates_text += f" | {', '.join(truncated['topics'])}"
            if truncated['notes']:
                candidates_text += f"\n    {truncated['notes']}"

        prompt = custom_prompt.format(
            new_entry=new_entry_text,
            candidates=candidates_text
        )

        # Call OpenAI Responses API with Structured Outputs
        request_data = json.dumps({
            'model': 'gpt-5-nano',
            'input': prompt,
            'text': {
                'format': {
                    'type': 'json_schema',
                    'name': 'content_similarity',
                    'strict': True,
                    'schema': {
                        'type': 'object',
                        'properties': {
                            'similar': {
                                'type': 'array',
                                'items': {
                                    'type': 'object',
                                    'properties': {
                                        'timestamp': {'type': 'string'},
                                        'score': {'type': 'integer'},
                                        'reason': {'type': 'string'}
                                    },
                                    'required': ['timestamp', 'score', 'reason'],
                                    'additionalProperties': False
                                }
                            }
                        },
                        'required': ['similar'],
                        'additionalProperties': False
                    }
                }
            }
        }).encode('utf-8')

        req = urllib.request.Request(
            'https://api.openai.com/v1/responses',
            data=request_data,
            headers={
                'Content-Type': 'application/json',
                'Authorization': f'Bearer {api_key}'
            }
        )

        with urllib.request.urlopen(req, timeout=60) as response:
            result = json.loads(response.read().decode('utf-8'))

        # Extract text from response - guaranteed valid JSON with structured outputs
        for item in result.get('output', []):
            if item.get('type') == 'message':
                for content_item in item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        text_content = content_item.get('text', '')
                        parsed = json.loads(text_content)
                        return parsed.get('similar', [])

        return []

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: find_similar_entries error: {str(e)}\n")
        return []


def add_relationships_to_entry(project_folder, timestamp, related_entries):
    """Add Related: metadata line to an existing kb.md entry."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return False

        folder_path = result
        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return False

        with open(kb_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # Find entry with matching timestamp
        i = 0
        while i < len(lines):
            if f'`{timestamp}`' in lines[i]:
                # Find insertion point (after Notes, before Entity/Topics/etc)
                j = i + 1
                insert_pos = j

                while j < len(lines):
                    line = lines[j]
                    if line.startswith('- `'):  # Next entry
                        break
                    if line.startswith('  - Notes:') or line.startswith('  - ParentId:'):
                        insert_pos = j + 1
                        # Skip multi-line notes
                        while insert_pos < len(lines) and lines[insert_pos].startswith('    '):
                            insert_pos += 1
                    elif line.startswith('  - Entity:') or line.startswith('  - Topics:'):
                        break
                    j += 1

                # Format related entries: timestamp (type), timestamp (type)
                related_str = ', '.join([f"{r['timestamp']} ({r['type']})" for r in related_entries])
                lines.insert(insert_pos, f"  - Related: {related_str}\n")

                with open(kb_file, 'w', encoding='utf-8') as f:
                    f.writelines(lines)

                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Added relationships to {timestamp}\n")
                return True

            i += 1

        return False

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: add_relationships_to_entry error: {str(e)}\n")
        return False


def add_similar_to_entry(project_folder, timestamp, similar_entries):
    """Add Similar: metadata line to an existing kb.md entry."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return False

        folder_path = result
        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return False

        with open(kb_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # Find entry with matching timestamp
        i = 0
        while i < len(lines):
            if f'`{timestamp}`' in lines[i]:
                # Find insertion point (after Notes/Related, before Entity/Topics)
                j = i + 1
                insert_pos = j

                while j < len(lines):
                    line = lines[j]
                    if line.startswith('- `'):  # Next entry
                        break
                    if line.startswith('  - Notes:') or line.startswith('  - ParentId:') or line.startswith('  - Related:'):
                        insert_pos = j + 1
                        # Skip multi-line notes
                        while insert_pos < len(lines) and lines[insert_pos].startswith('    '):
                            insert_pos += 1
                    elif line.startswith('  - Entity:') or line.startswith('  - Topics:'):
                        break
                    j += 1

                # Format similar entries: timestamp (score), timestamp (score)
                similar_str = ', '.join([f"{s['timestamp']} ({s['score']})" for s in similar_entries])
                lines.insert(insert_pos, f"  - Similar: {similar_str}\n")

                with open(kb_file, 'w', encoding='utf-8') as f:
                    f.writelines(lines)

                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Added similar entries to {timestamp}\n")
                return True

            i += 1

        return False

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: add_similar_to_entry error: {str(e)}\n")
        return False


def background_process_entry(project_folder, timestamp, entry, api_key, file_path=None):
    """
    Background thread for AI processing pipeline.

    Runs after the initial save returns to the UI, performing:
    1. AI summary generation (type-specific)
    2. Classification (entity, topics, people, work/personal, grammar correction)
    3. Relationship extraction (parse notes for references to other entries)
    4. Content similarity (find semantically similar entries)

    Each step's results are written back to kb.md incrementally.
    Prompts are loaded from settings.json file.

    Args:
        project_folder (str): Path to project folder containing kb.md.
        timestamp (str): Entry timestamp for identification.
        entry (dict): Original entry data.
        api_key (str): OpenAI API key.
        file_path (str|None): Path to associated file for summarisation.
    """
    log_file = Path(__file__).parent / 'host.log'

    try:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: Background thread started for {timestamp}\n")

        # Load prompts from settings.json (single source of truth)
        settings = load_settings()
        classification_prompt = get_prompt(settings, 'classification_prompt')

        # Load all summary prompts
        summary_prompts = {
            'image': get_prompt(settings, 'image_prompt'),
            'audio': get_prompt(settings, 'audio_prompt'),
            'document': get_prompt(settings, 'document_prompt'),
            'link': get_prompt(settings, 'link_prompt'),
            'text': get_prompt(settings, 'text_prompt'),
            'research': get_prompt(settings, 'research_prompt'),
        }

        # STEP 1: AI Summary (before classification so summary can inform classification)
        summary = None
        if api_key:
            entry_type = entry.get('type', '')
            summary = generate_ai_summary(entry_type, entry, file_path, api_key, summary_prompts)
            if summary:
                add_summary_to_entry(project_folder, timestamp, summary)

        # STEP 2: Classification (includes grammar correction, entity, topics, people, work/personal)
        classification = None
        if api_key:
            existing_topics = load_topics(project_folder)
            existing_entities = load_entities(project_folder)
            classification = classify_entry(entry, api_key, existing_topics, existing_entities, summary, classification_prompt)

            if classification:
                # Apply grammar-corrected notes
                corrected_notes = classification.get('corrected_notes', '').strip()
                original_notes = entry.get('notes', '').strip()
                if corrected_notes and corrected_notes != original_notes:
                    update_entry_notes(project_folder, timestamp, corrected_notes)
                    entry['notes'] = corrected_notes

                # Save topics/entities to JSON files
                new_topics = classification.get('topics', [])
                new_entities = classification.get('people', [])

                if new_topics:
                    save_topics(project_folder, existing_topics + new_topics)
                if new_entities:
                    save_entities(project_folder, existing_entities + new_entities)

                # Add classification metadata to kb.md entry
                add_classification_to_entry(project_folder, timestamp, classification)

                # Add work/personal category
                category = classification.get('category')
                if category:
                    add_category_to_entry(project_folder, timestamp, category)

        # STEP 3: Relationship extraction (parse notes for references to other entries)
        relationship_extraction_prompt = get_prompt(settings, 'relationship_extraction_prompt')
        relationship_resolution_prompt = get_prompt(settings, 'relationship_resolution_prompt')

        notes = entry.get('notes', '').strip()
        if notes and len(notes) >= 10 and api_key:
            # Update entry with classification info for extraction
            entry['topics'] = classification.get('topics', []) if classification else []

            relationships = extract_relationships(entry, api_key, relationship_extraction_prompt)
            if relationships.get('has_relationships') and relationships.get('references'):
                # Load all entries for candidate search
                all_entries = load_all_entries(project_folder)

                resolved_relationships = []
                for ref in relationships['references']:
                    # Find candidates using fuzzy search
                    candidates = fuzzy_search_candidates(
                        all_entries,
                        ref.get('target', ''),
                        ref.get('type', 'any'),
                        ref.get('temporal', 'none'),
                        timestamp
                    )

                    if candidates:
                        # Use AI to pick the best match
                        resolution = resolve_with_ai(
                            ref.get('phrase', ''),
                            ref.get('target', ''),
                            candidates,
                            api_key,
                            relationship_resolution_prompt
                        )

                        if resolution.get('timestamp') and resolution.get('confidence') in ['high', 'medium']:
                            resolved_relationships.append({
                                'timestamp': resolution['timestamp'],
                                'type': 'mentions'
                            })

                if resolved_relationships:
                    add_relationships_to_entry(project_folder, timestamp, resolved_relationships)

        # STEP 4: Content similarity (find semantically similar entries)
        content_similarity_prompt = get_prompt(settings, 'content_similarity_prompt')

        # Check if entry has meaningful content for comparison
        has_content = (
            entry.get('notes', '').strip() or
            entry.get('selectedText', '').strip() or
            summary
        )

        if has_content and api_key:
            # Update entry with all available context
            entry['aiSummary'] = summary or ''
            entry['topics'] = classification.get('topics', []) if classification else []
            entry['entity'] = classification.get('entity', '') if classification else ''
            entry['people'] = classification.get('people', []) if classification else []

            # Load all entries if not already loaded
            if 'all_entries' not in locals():
                all_entries = load_all_entries(project_folder)

            # Get pre-filtered candidates
            similarity_candidates = get_similarity_candidates(all_entries, entry)

            if similarity_candidates:
                similar = find_similar_entries(entry, similarity_candidates, api_key, content_similarity_prompt)
                if similar:
                    add_similar_to_entry(project_folder, timestamp, similar)

        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: Background thread completed for {timestamp}\n")

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: Background thread error: {str(e)}\n")


def add_classification_to_entry(project_folder, timestamp, classification):
    """Add Entity/Topics/People metadata lines to an existing kb.md entry."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return False

        folder_path = result
        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return False

        with open(kb_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # Find entry with matching timestamp
        i = 0
        while i < len(lines):
            if f'`{timestamp}`' in lines[i]:
                # Found entry - find where to insert classification
                j = i + 1
                insert_pos = j

                while j < len(lines):
                    if lines[j].startswith('- '):
                        # Hit next entry
                        insert_pos = j
                        break
                    elif lines[j].strip() == '':
                        # Hit blank line (end of entry)
                        insert_pos = j
                        break
                    j += 1
                else:
                    insert_pos = len(lines)

                # Build classification lines
                class_lines = []
                entity = classification.get('entity', '')
                topics = classification.get('topics', [])
                people = classification.get('people', [])

                if entity:
                    class_lines.append(f"  - Entity: {entity}\n")
                    # If classified as task, add default Status for kanban board
                    if entity == 'task':
                        class_lines.append(f"  - Status: not-started\n")
                if topics:
                    class_lines.append(f"  - Topics: {', '.join(topics)}\n")
                if people:
                    class_lines.append(f"  - People: {', '.join(people)}\n")

                # Insert classification lines
                for idx, cl in enumerate(class_lines):
                    lines.insert(insert_pos + idx, cl)

                # Write back
                with open(kb_file, 'w', encoding='utf-8') as f:
                    f.writelines(lines)

                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Classification added for {timestamp}: {classification}\n")

                return True

            i += 1

        return False

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: add_classification_to_entry error: {str(e)}\n")
        return False


def add_category_to_entry(project_folder, timestamp, category):
    """Add Category (work/personal) metadata line to an existing kb.md entry."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return False

        folder_path = result
        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return False

        with open(kb_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # Find entry with matching timestamp
        i = 0
        while i < len(lines):
            if f'`{timestamp}`' in lines[i]:
                # Found entry - find where to insert category
                j = i + 1
                insert_pos = j

                while j < len(lines):
                    if lines[j].startswith('- '):
                        # Hit next entry
                        insert_pos = j
                        break
                    elif lines[j].strip() == '':
                        # Hit blank line (end of entry)
                        insert_pos = j
                        break
                    j += 1
                else:
                    insert_pos = len(lines)

                # Insert category line
                category_line = f"  - Category: {category}\n"
                lines.insert(insert_pos, category_line)

                # Write back
                with open(kb_file, 'w', encoding='utf-8') as f:
                    f.writelines(lines)

                with open(log_file, 'a', encoding='utf-8') as f:
                    f.write(f"{datetime.now()}: Category added for {timestamp}: {category}\n")

                return True

            i += 1

        return False

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: add_category_to_entry error: {str(e)}\n")
        return False


def classify_and_update_entry(project_folder, timestamp, entry, api_key):
    """Classify entry and update kb.md with Entity/Topics/People metadata (background process)."""
    log_file = Path(__file__).parent / 'host.log'

    try:
        # Validate project folder
        is_valid, result = validate_project_folder(project_folder)
        if not is_valid:
            return {'success': False, 'error': result}

        folder_path = result
        kb_file = folder_path / 'kb.md'

        if not kb_file.exists():
            return {'success': False, 'error': 'kb.md file does not exist'}

        # Load existing topics and entities
        existing_topics = load_topics(project_folder)
        existing_entities = load_entities(project_folder)

        # Perform AI classification
        classification = classify_entry(entry, api_key, existing_topics, existing_entities)

        if not classification:
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Classification returned no result for {timestamp}\n")
            return {'success': True, 'message': 'No classification data'}

        # Update topics and entities JSON files
        new_topics = classification.get('topics', [])
        new_entities = classification.get('people', [])

        if new_topics:
            save_topics(project_folder, existing_topics + new_topics)
        if new_entities:
            save_entities(project_folder, existing_entities + new_entities)

        # Read kb.md and find entry by timestamp
        with open(kb_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # Find the entry with matching timestamp and add classification metadata
        updated = False
        i = 0
        while i < len(lines):
            line = lines[i]

            if f'`{timestamp}`' in line:
                # Found the entry - find where to insert classification metadata
                # Look for end of this entry (next main entry line or end of file)
                j = i + 1
                insert_pos = j

                while j < len(lines):
                    if lines[j].startswith('- '):
                        # Hit next entry
                        insert_pos = j
                        break
                    elif lines[j].strip() == '':
                        # Hit blank line (end of entry)
                        insert_pos = j
                        break
                    j += 1
                else:
                    insert_pos = len(lines)

                # Build classification lines
                class_lines = []
                entity = classification.get('entity', '')
                topics = classification.get('topics', [])
                people = classification.get('people', [])

                if entity:
                    class_lines.append(f"  - Entity: {entity}\n")
                if topics:
                    class_lines.append(f"  - Topics: {', '.join(topics)}\n")
                if people:
                    class_lines.append(f"  - People: {', '.join(people)}\n")

                # Insert classification lines before blank line / next entry
                for idx, cl in enumerate(class_lines):
                    lines.insert(insert_pos + idx, cl)

                updated = True
                break

            i += 1

        if updated:
            with open(kb_file, 'w', encoding='utf-8') as f:
                f.writelines(lines)
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Classification metadata added for {timestamp}\n")
            return {'success': True, 'message': 'Classification added', 'classification': classification}
        else:
            return {'success': False, 'error': f'Entry with timestamp {timestamp} not found'}

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: classify_and_update_entry error: {str(e)}\n")
        return {'success': False, 'error': str(e)}


def launch_widget():
    """Launch widget with single-instance check."""
    import subprocess
    import tempfile

    log_file = Path(__file__).parent / 'host.log'

    def log(msg):
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: [launch_widget] {msg}\n")

    log("=" * 50)
    log("launch_widget() called")

    widget_path = Path(__file__).parent / 'widget.pyw'
    state_file = Path(__file__).parent / 'widget_state.json'
    lock_file = Path(tempfile.gettempdir()) / 'ultrathink_widget.lock'

    log(f"widget_path: {widget_path}")
    log(f"widget_path.exists(): {widget_path.exists()}")

    if not widget_path.exists():
        log("ERROR: Widget not found!")
        return {'success': False, 'error': 'Widget not found'}

    # Check state file for running widget
    log(f"Checking state file: {state_file}")
    log(f"state_file.exists(): {state_file.exists()}")
    try:
        if state_file.exists():
            with open(state_file, 'r') as f:
                state = json.load(f)
            log(f"State file contents: {state}")
            if state.get('running') and state.get('pid'):
                log(f"State says running with PID {state['pid']} - checking...")
                result = subprocess.run(
                    ['tasklist', '/FI', f'PID eq {state["pid"]}', '/NH'],
                    capture_output=True, text=True
                )
                log(f"tasklist result: {result.stdout.strip()}")
                if 'pythonw.exe' in result.stdout or 'python.exe' in result.stdout:
                    log("Process IS running - trying to focus window")
                    # Focus existing window
                    try:
                        import ctypes
                        hwnd = ctypes.windll.user32.FindWindowW(None, "UltraThink Widget")
                        log(f"FindWindowW returned: {hwnd}")
                        if hwnd:
                            ctypes.windll.user32.SetForegroundWindow(hwnd)
                            log("Focused existing window")
                    except Exception as e:
                        log(f"Focus error: {e}")
                    return {'success': True, 'already_running': True}
                else:
                    log("Process NOT running - state is stale")
    except Exception as e:
        log(f"State check error: {e}")

    # Clean stale lock file
    log(f"Cleaning stale lock file: {lock_file}")
    try:
        lock_file.unlink(missing_ok=True)
        log("Lock file cleaned")
    except Exception as e:
        log(f"Lock cleanup error: {e}")

    log("Launching widget process...")
    try:
        # Launch with pythonw (no console window)
        process = subprocess.Popen(
            ['pythonw', str(widget_path)],
            creationflags=subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP,
            close_fds=True
        )
        log(f"SUCCESS: Widget launched with PID {process.pid}")
        return {'success': True, 'pid': process.pid}
    except FileNotFoundError:
        log("pythonw not found - trying python instead")
        # Try with python if pythonw not found
        try:
            process = subprocess.Popen(
                ['python', str(widget_path)],
                creationflags=subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP,
                close_fds=True
            )
            log(f"SUCCESS: Widget launched with python, PID {process.pid}")
            return {'success': True, 'pid': process.pid}
        except Exception as e:
            log(f"ERROR: Failed to launch with python: {e}")
            return {'success': False, 'error': f'Failed to launch: {str(e)}'}
    except Exception as e:
        log(f"ERROR: Launch failed: {e}")
        return {'success': False, 'error': str(e)}


def close_widget():
    """Close widget by PID from state file."""
    import subprocess
    import tempfile

    state_file = Path(__file__).parent / 'widget_state.json'
    lock_file = Path(tempfile.gettempdir()) / 'ultrathink_widget.lock'
    pid = None

    # Get PID from state file
    try:
        if state_file.exists():
            with open(state_file, 'r') as f:
                pid = json.load(f).get('pid')
    except Exception:
        pass

    try:
        kill_succeeded = False
        if pid:
            # Kill by specific PID - precise and safe
            result = subprocess.run(['taskkill', '/F', '/PID', str(pid)], capture_output=True)
            kill_succeeded = result.returncode == 0
        else:
            # Fallback: exact window title match (more targeted than wildcard)
            result = subprocess.run([
                'powershell', '-Command',
                "Get-Process | Where-Object {$_.MainWindowTitle -eq 'UltraThink Widget'} | Stop-Process -Force"
            ], capture_output=True)
            # PowerShell returns 0 even if no process found, so consider it success
            kill_succeeded = result.returncode == 0

        # Cleanup state file (always do this, even if kill "failed" - widget may have already exited)
        try:
            with open(state_file, 'w') as f:
                json.dump({'running': False, 'pid': None}, f)
        except Exception:
            pass

        # Cleanup lock file
        try:
            lock_file.unlink(missing_ok=True)
        except Exception:
            pass

        return {'success': True}
    except Exception as e:
        return {'success': False, 'error': str(e)}


def browse_folder():
    """Open folder picker dialog and return selected path."""
    try:
        import tkinter as tk
        from tkinter import filedialog

        root = tk.Tk()
        root.withdraw()  # Hide main window
        root.attributes('-topmost', True)  # Bring dialog to front

        folder = filedialog.askdirectory(title="Select UltraThink Project Folder")
        root.destroy()

        if folder:
            # Normalize to Windows path with trailing backslash
            folder = folder.replace('/', '\\')
            if not folder.endswith('\\'):
                folder += '\\'
            return {'success': True, 'path': folder}
        return {'success': False, 'error': 'No folder selected'}
    except Exception as e:
        return {'success': False, 'error': str(e)}


def main():
    """Main loop for native messaging."""
    # Log to file for debugging (optional)
    log_file = Path(__file__).parent / 'host.log'

    try:
        while True:
            # Read message from extension
            message = read_message()

            if message is None:
                break

            # Debug logging
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write(f"{datetime.now()}: Received message: {json.dumps(message)}\n")

            # Handle the message
            if message.get('action') == 'append':
                project_folder = message.get('projectFolder')
                entry = message.get('entry')
                api_key = message.get('apiKey')  # API key for AI classification
                # Note: prompts are now read from settings.json by background_process_entry

                # Persist projectFolder to settings.json so widget/kb-server stay in sync
                if project_folder:
                    current_settings = load_settings()
                    if current_settings.get('project_folder') != project_folder:
                        update_settings({'project_folder': project_folder})

                result = append_to_kb(project_folder, entry, api_key)

                # Extract background task before sending (don't send internal field to extension)
                background_task = result.pop('_background_task', None)

                # Send response immediately so UI is unblocked
                send_message(result)

                # Now run background processing (grammar + classification + summary)
                # Prompts are read from settings.json by background_process_entry
                if background_task:
                    try:
                        background_process_entry(
                            background_task['project_folder'],
                            background_task['timestamp'],
                            background_task['entry'],
                            background_task['api_key'],
                            background_task.get('file_path')
                        )
                    except Exception as bg_error:
                        with open(log_file, 'a', encoding='utf-8') as f:
                            f.write(f"{datetime.now()}: Background processing error: {str(bg_error)}\n")

                continue  # Skip the send_message at the end since we already sent
            elif message.get('action') == 'update_last_entry':
                project_folder = message.get('projectFolder')
                timestamp = message.get('timestamp')
                new_content = message.get('newContent')

                result = update_last_entry(project_folder, timestamp, new_content)
                send_message(result)
            elif message.get('action') == 'launch_widget':
                result = launch_widget()
                send_message(result)
            elif message.get('action') == 'close_widget':
                result = close_widget()
                send_message(result)
            elif message.get('action') == 'browse_folder':
                result = browse_folder()
                send_message(result)
            elif message.get('action') == 'classify_entry':
                # Background classification - called after grammar fix completes
                project_folder = message.get('projectFolder')
                timestamp = message.get('timestamp')
                entry = message.get('entry')
                api_key = message.get('apiKey')

                result = classify_and_update_entry(project_folder, timestamp, entry, api_key)
                send_message(result)
            elif message.get('action') == 'search_github':
                # Search GitHub for issues/commits related to a query
                query = message.get('query')
                github_token = message.get('githubToken')
                repos = message.get('githubRepos')
                max_results = message.get('maxResults', 10)

                result = search_github(query, github_token, repos, max_results)
                send_message(result)
            elif message.get('action') == 'get_settings':
                # Return settings (for extension options page)
                settings = load_settings()
                send_message({
                    'success': True,
                    'project_folder': settings.get('project_folder', ''),
                    'openai_api_key': settings.get('openai_api_key', ''),
                    'github_token': settings.get('github_token', ''),
                    'github_org': settings.get('github_org', ''),
                    'github_repos': settings.get('github_repos', ''),
                    'notion_token': settings.get('notion_token', ''),
                    'fastmail_token': settings.get('fastmail_token', ''),
                    'capsule_token': settings.get('capsule_token', ''),
                    'extension_id': settings.get('extension_id', '')
                })
            elif message.get('action') == 'save_settings':
                # Save settings from extension (settings.json is source of truth)
                updates = message.get('settings', {})
                if update_settings(updates):
                    send_message({'success': True})
                else:
                    send_message({'success': False, 'error': 'Failed to save settings'})
            else:
                send_message({'success': False, 'error': 'Unknown action'})

    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now()}: Error: {str(e)}\n")
        send_message({'success': False, 'error': str(e)})


if __name__ == '__main__':
    main()
